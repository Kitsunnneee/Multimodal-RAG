"""Document processing utilities for the Multimodal RAG system."""
import os
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union, Any

import pandas as pd
from langchain_core.documents import Document
from langchain_text_splitters import CharacterTextSplitter
from pptx import Presentation
from unstructured.partition.pdf import partition_pdf

from .llama_parse_loader import LlamaParseLoader
from .file_utils import get_file_type, is_supported_file
from langchain_community.document_loaders import (
    UnstructuredPowerPointLoader,
    UnstructuredExcelLoader,
    CSVLoader,
    TextLoader,
    UnstructuredFileLoader
)

from .config import (
    CHUNK_SIZE,
    CHUNK_OVERLAP,
    MAX_CHARACTERS,
    NEW_AFTER_N_CHARS,
    COMBINE_TEXT_UNDER_N_CHARS,
    DATA_DIR,
)
from .utils import encode_image, is_image_data, split_image_text_types


class DocumentProcessor:
    """Process documents and extract text, tables, and images."""
    
    def __init__(
        self, 
        output_dir: Optional[Union[str, Path]] = None,
        use_llama_parse: bool = True,
        llama_parse_kwargs: Optional[Dict[str, Any]] = None,
    ):
        """Initialize the document processor.
        
        Args:
            output_dir: Directory to save extracted images. If None, uses DATA_DIR.
            use_llama_parse: Whether to use LlamaParse for document parsing
            llama_parse_kwargs: Additional arguments to pass to LlamaParse
        """
        self.output_dir = Path(output_dir) if output_dir else DATA_DIR
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Initialize LlamaParse if enabled and API key is available
        self.use_llama_parse = use_llama_parse
        self.llama_loader = None
        
        if use_llama_parse:
            try:
                self.llama_loader = LlamaParseLoader(
                    output_dir=self.output_dir,
                    **(llama_parse_kwargs or {})
                )
                print("LlamaParse initialized successfully")
            except Exception as e:
                print(f"Warning: Failed to initialize LlamaParse: {e}")
                print("Falling back to default parser")
                self.use_llama_parse = False
    
    def process_file(
        self,
        file_path: Union[str, Path],
        extract_images: bool = True,
        infer_table_structure: bool = True,
        use_llama_parse: Optional[bool] = None,
        **kwargs,
    ) -> Dict[str, List[Document]]:
        """Process a file (PDF or image) and extract text and images.
        
        Args:
            file_path: Path to the file (PDF or image)
            extract_images: Whether to extract images from the file (for PDFs)
            infer_table_structure: Whether to infer table structure (for PDFs)
            use_llama_parse: Override for using LlamaParse. If None, uses class default.
            **kwargs: Additional arguments to pass to the loader
            
        Returns:
            Dictionary containing 'texts', 'tables', and 'images' as lists of Documents
            
        Raises:
            FileNotFoundError: If the input file doesn't exist
            ValueError: If the file type is not supported
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_type = file_path.suffix.lower()  # Keep the dot for consistency
        
        # Initialize result dictionary
        result = {
            "texts": [],
            "tables": [],
            "images": []
        }
        
        # Use LlamaParse if enabled and available for this file type
        use_llama = use_llama_parse if use_llama_parse is not None else self.use_llama_parse
        if use_llama and self.llama_loader and file_type in ['.pdf', '.docx', '.pptx', '.xlsx']:
            try:
                # Remove the dot for the file_type parameter since LlamaParse expects it without dot
                documents = self.llama_loader.load_file(file_path, file_type=file_type.lstrip('.'), **kwargs)
                result["texts"] = documents
                return result
            except Exception as e:
                print(f"Warning: LlamaParse failed with error: {e}. Falling back to default parser.")
                if use_llama_parse:  # Only if explicitly requested, otherwise just warn
                    raise
        
        # Process the file based on its type
        if file_type == '.pdf':
            return self._process_pdf(
                file_path, 
                extract_images=kwargs.get('extract_images', True),
                infer_table_structure=kwargs.get('infer_table_structure', True)
            )
        elif file_type in ['.pptx', '.ppt']:
            return self._process_pptx(file_path)
        elif file_type in ['.xlsx', '.xls']:
            return self._process_excel(file_path)
        elif file_type == '.csv':
            return self._process_csv(file_path)
        elif file_type in ['.txt', '.md', '.markdown']:
            return self._process_text(file_path)
        elif file_type in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
            return self._process_image(file_path)
        else:
            # Try with unstructured loader as fallback
            try:
                loader = UnstructuredFileLoader(str(file_path))
                docs = loader.load()
                return {"texts": docs, "tables": [], "images": []}
            except Exception as e:
                raise ValueError(
                    f"Unsupported file type: {file_type}. "
                    "Supported types: .pdf, .pptx, .ppt, .xlsx, .xls, .csv, .txt, "
                    ".png, .jpg, .jpeg, .tiff, .bmp, .gif"
                )
    
    def _process_image(self, file_path: Path) -> Dict[str, List[Document]]:
        """Process an image file and extract text.
        
        Args:
            file_path: Path to the image file
            
        Returns:
            Dictionary containing 'texts' and empty 'tables' and 'images' lists
        """
        from .utils import encode_image, extract_text_from_image
        import os
        
        try:
            # Encode image to base64
            b64_image = encode_image(file_path)
            
            # Get the service account path from environment or use default location
            service_account_path = os.getenv('GOOGLE_APPLICATION_CREDENTIALS')
            if not service_account_path or not os.path.exists(service_account_path):
                service_account_path = os.path.join(
                    os.path.dirname(os.path.dirname(__file__)), 
                    'elite-thunder-461308-f7-cc85c56bb209.json'
                )
            
            if not os.path.exists(service_account_path):
                raise FileNotFoundError(
                    f"Service account file not found at {service_account_path}. "
                    "Please set the GOOGLE_APPLICATION_CREDENTIALS environment variable "
                    "to point to your service account JSON file."
                )
            
            # Extract text using OCR with Gemini API
            extracted_text = extract_text_from_image(
                b64_image,
                service_account_path=service_account_path
            )
            
            if not extracted_text:
                print(f"Warning: No text could be extracted from {file_path}")
                return {"texts": [], "tables": [], "images": []}
            
            # Store the base64-encoded image in the document content
            # along with the extracted text for better retrieval
            doc = Document(
                page_content=f"[IMAGE] {extracted_text}",
                metadata={
                    "source": str(file_path),
                    "page": 0,
                    "type": "image",
                    "original_content": extracted_text,
                    "image_data": b64_image  # Store the base64-encoded image data
                }
            )
            
            return {"texts": [doc], "tables": [], "images": [doc]}
            
        except Exception as e:
            print(f"Error processing image {file_path}: {e}")
            # Return empty results on error
            return {"texts": [], "tables": [], "images": []}
    
    def _process_pptx(self, file_path: Path) -> Dict[str, List[Document]]:
        """Process a PowerPoint file and extract text and slides.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary containing 'texts' and empty 'tables' and 'images' lists
        """
        try:
            loader = UnstructuredPowerPointLoader(str(file_path))
            docs = loader.load()
            return {"texts": docs, "tables": [], "images": []}
        except Exception as e:
            print(f"Error processing PowerPoint file {file_path}: {e}")
            return {"texts": [], "tables": [], "images": []}
    
    def _process_excel(self, file_path: Path) -> Dict[str, List[Document]]:
        """Process an Excel file and extract data from all sheets.
        
        Args:
            file_path: Path to the Excel file
            
        Returns:
            Dictionary containing 'texts' (for sheet data) and 'tables' (for structured data)
        """
        try:
            # First try with UnstructuredExcelLoader
            loader = UnstructuredExcelLoader(str(file_path), mode="elements")
            docs = loader.load()
            
            # Separate text and tables
            texts = []
            tables = []
            
            for doc in docs:
                if doc.metadata.get("category") == "Table":
                    tables.append(doc)
                else:
                    texts.append(doc)
            
            return {"texts": texts, "tables": tables, "images": []}
            
        except Exception as e:
            print(f"Error processing Excel file {file_path}: {e}")
            return {"texts": [], "tables": [], "images": []}
    
    def _process_csv(self, file_path: Path) -> Dict[str, List[Document]]:
        """Process a CSV file and extract data.
        
        Args:
            file_path: Path to the CSV file
            
        Returns:
            Dictionary containing 'texts' with CSV data
        """
        try:
            # First try with pandas for better handling of various CSV formats
            try:
                import pandas as pd
                df = pd.read_csv(file_path, nrows=1000)  # Read first 1000 rows
                text = df.to_string()
                doc = Document(
                    page_content=text,
                    metadata={"source": str(file_path), "type": "csv"}
                )
                return {"texts": [doc], "tables": [], "images": []}
            except Exception as pd_error:
                # Fall back to CSVLoader if pandas fails
                loader = CSVLoader(str(file_path))
                docs = loader.load()
                return {"texts": docs, "tables": [], "images": []}
                
        except Exception as e:
            print(f"Error processing CSV file {file_path}: {e}")
            return {"texts": [], "tables": [], "images": []}
    
    def _process_text(self, file_path: Path) -> Dict[str, List[Document]]:
        """Process a plain text file.
        
        Args:
            file_path: Path to the text file
            
        Returns:
            Dictionary containing 'texts' with file content
        """
        try:
            loader = TextLoader(str(file_path))
            docs = loader.load()
            return {"texts": docs, "tables": [], "images": []}
        except Exception as e:
            print(f"Error processing text file {file_path}: {e}")
            return {"texts": [], "tables": [], "images": []}
    
    def _process_pdf(
        self,
        file_path: Path,
        extract_images: bool = True,
        infer_table_structure: bool = True,
    ) -> Dict[str, List[Document]]:
        """Process a PDF file and extract text, tables, and images."""
        from unstructured.partition.pdf import partition_pdf
        from .config import MAX_CHARACTERS, NEW_AFTER_N_CHARS, COMBINE_TEXT_UNDER_N_CHARS
        
        try:
            # Ensure output directory exists
            self.output_dir.mkdir(parents=True, exist_ok=True)
            
            # Extract elements from PDF with error handling
            try:
                elements = partition_pdf(
                    filename=str(file_path),
                    extract_images_in_pdf=extract_images,
                    infer_table_structure=infer_table_structure,
                    chunking_strategy="by_title",
                    max_characters=MAX_CHARACTERS,
                    new_after_n_chars=NEW_AFTER_N_CHARS,
                    combine_text_under_n_chars=COMBINE_TEXT_UNDER_N_CHARS,
                    image_output_dir_path=str(self.output_dir),
                )
            except Exception as e:
                raise ValueError(f"Failed to process PDF: {str(e)}") from e
            
            # Categorize elements with type checking
            texts = []
            tables = []
            
            for element in elements:
                try:
                    element_str = str(element or '').strip()
                    if not element_str:
                        continue
                        
                    element_type = str(type(element))
                    if "Table" in element_type:
                        tables.append(element_str)
                    elif any(t in element_type for t in ["CompositeElement", "Text"]):
                        texts.append(element_str)
                except Exception as e:
                    print(f"Warning: Error processing element: {e}")
                    continue
            
            # Convert to Documents with content validation
            text_docs = [Document(page_content=text) for text in texts if text and len(text) > 10]
            table_docs = [Document(page_content=table) for table in tables if table]
            
            # Process images if any were extracted
            image_docs = []
            if extract_images:
                image_files = list(self.output_dir.glob("*.png")) + list(self.output_dir.glob("*.jpg"))
                for img_path in image_files:
                    try:
                        if img_path.stat().st_size == 0:
                            continue
                        b64_image = encode_image(img_path)
                        if b64_image and len(b64_image) > 100:  # Basic validation
                            image_docs.append(Document(
                                page_content=b64_image, 
                                metadata={"source": str(img_path.name)}
                            ))
                        # Clean up the image file after processing
                        img_path.unlink(missing_ok=True)
                    except Exception as e:
                        print(f"Warning: Error processing image {img_path}: {e}")
                        continue
            
            return {
                "texts": text_docs,
                "tables": table_docs,
                "images": image_docs,
            }
            
        except Exception as e:
            # Clean up any temporary files on error
            if 'elements' in locals():
                del elements
            raise Exception(f"Error processing PDF {file_path.name}: {str(e)}") from e
    
    def _process_pptx(self, file_path: Union[str, Path], extra_info: Optional[Dict[str, Any]] = None) -> List[Document]:
        """Extract text from a PowerPoint presentation.
        
        Args:
            file_path: Path to the PowerPoint file
            extra_info: Additional metadata to include in the documents
            
        Returns:
            List of Document objects, one per slide
        """
        prs = Presentation(file_path)
        documents = []
        
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            
            if text:
                content = "\n".join(filter(None, text))
                metadata = {
                    "source": str(file_path),
                    "page": i + 1,
                    "file_type": "pptx",
                    "total_slides": len(prs.slides),
                    "slide_number": i + 1,
                    **(extra_info or {})
                }
                documents.append(Document(page_content=content, metadata=metadata))
        
        return documents
    
    def _process_xlsx(self, file_path: Union[str, Path], extra_info: Optional[Dict[str, Any]] = None) -> List[Document]:
        """Extract data from an Excel file.
        
        Args:
            file_path: Path to the Excel file
            extra_info: Additional metadata to include in the documents
            
        Returns:
            List of Document objects, one per sheet
        """
        xls = pd.ExcelFile(file_path)
        documents = []
        
        for sheet_name in xls.sheet_names:
            try:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                
                # Convert DataFrame to markdown for better formatting
                content = f"# {sheet_name}\n\n{df.to_markdown()}"
                
                metadata = {
                    "source": str(file_path),
                    "file_type": "xlsx",
                    "sheet_name": sheet_name,
                    "shape": f"{df.shape[0]} rows x {df.shape[1]} columns",
                    **(extra_info or {})
                }
                documents.append(Document(page_content=content, metadata=metadata))
            except Exception as e:
                print(f"Error processing sheet '{sheet_name}': {e}")
        
        return documents

    def _process_file(
        self, 
        file_path: Union[str, Path],
        file_type: Optional[str] = None,
        extra_info: Optional[Dict[str, Any]] = None,
    ) -> List[Document]:
        """Process a single file and return a list of documents.
        
        Args:
            file_path: Path to the file to process
            file_type: Type of the file (pdf, docx, pptx, xlsx, etc.)
            extra_info: Additional metadata to include in the documents
            
        Returns:
            List of Document objects
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        # Detect file type if not provided
        if not file_type:
            file_type, _ = get_file_type(file_path)
        
        # Use LlamaParse if available and not disabled
        if self.use_llama_parse and self.llama_loader and file_type in ('pdf', 'docx'):
            try:
                return self.llama_loader.load_file(
                    file_path, 
                    file_type=file_type,
                    extra_info=extra_info
                )
            except Exception as e:
                print(f"Error using LlamaParse: {e}. Falling back to local processing.")
        
        # Process based on file type
        if not is_supported_file(file_path):
            raise ValueError(f"Unsupported file type: {file_type}")
            
        if file_type == 'pdf':
            return self._process_pdf(file_path, extra_info=extra_info)
        elif file_type == 'docx':
            return self._process_docx(file_path, extra_info=extra_info)
        elif file_type == 'pptx':
            return self._process_pptx(file_path, extra_info=extra_info)
        elif file_type == 'xlsx':
            return self._process_xlsx(file_path, extra_info=extra_info)
        else:
            # Default to text processing
            return self._process_text_file(file_path, extra_info=extra_info)


def generate_summaries(
    texts: List[str],
    model_name: str = "gemini-2.0-flash",
    temperature: float = 0.0,
    max_tokens: int = 1000,
) -> List[str]:
    """Generate summaries for a list of text chunks.
    
    Args:
        texts: List of text chunks to summarize
        model_name: Name of the model to use for summarization
        temperature: Temperature for text generation
        max_tokens: Maximum number of tokens to generate
        
    Returns:
        List of summaries
    """
    from langchain_core.prompts import PromptTemplate
    from langchain_core.output_parsers import StrOutputParser
    
    if not texts:
        return []
    
    prompt_template = """You are an assistant tasked with summarizing text for retrieval. 
    These summaries will be embedded and used to retrieve the raw text. 
    Give a concise summary of the text that is well optimized for retrieval.
    
    Text: {text}
    """
    
    prompt = PromptTemplate.from_template(prompt_template)
    
    try:
        # First try using VertexAI
        from langchain_google_vertexai import VertexAI
        
        llm = VertexAI(
            model_name=model_name,
            temperature=temperature,
            max_output_tokens=max_tokens,
        )
        chain = prompt | llm
        return chain.batch([{"text": text} for text in texts], {"max_concurrency": 5})
        
    except Exception as e:
        print(f"Warning: Could not use VertexAI: {e}")
        print("Falling back to a simple text truncation for summarization...")
        
        # Fallback: Just return the first few sentences as a simple summary
        def simple_summary(text: str, max_sentences: int = 3) -> str:
            import re
            # Split into sentences (very basic)
            sentences = re.split(r'(?<=[.!?])\s+', text)
            return ' '.join(sentences[:max_sentences])
            
        return [simple_summary(text) for text in texts]
